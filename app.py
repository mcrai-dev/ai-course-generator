import os
import re
import json
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

def clean_text(text):
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)         # Lignes non terminées
    text = re.sub(r'\s{2,}', ' ', text)                  # Espaces multiples
    text = re.sub(r'(\w)-\s+(\w)', r'\1\2', text)         # Coupures avec "-"
    text = re.sub(r'\[\s*\]', '', text)                  # [ ] orphelins
    return text.strip()

def is_code_line(text):
    """Détecte si une ligne ressemble à du code (tous langages)"""
    # Patterns génériques pour tous les langages
    code_patterns = [
        # Syntaxe générale
        r'^\s*[a-zA-Z_][a-zA-Z0-9_]*\s*\([^)]*\)\s*[\{:]',  # function() { ou function():
        r'^\s*[{}]\s*$',                                      # Accolades seules
        r'^\s*(if|else|while|for|do|switch|case|try|catch|finally)\s*[\(\{]',  # Structures de contrôle
        r'^\s*(def|function|func|sub|proc|void|int|string|bool|float|double|char|var|let|const)\s+\w+',  # Déclarations
        r'[=!<>]=|[+\-*/]=|\+\+|--|&&|\|\||<<|>>',           # Opérateurs
        r'^\s*[a-zA-Z_][a-zA-Z0-9_]*\s*=[^=]',               # Assignations
        r'^\s*[a-zA-Z_][a-zA-Z0-9_]*\.\w+\s*[\(=]',          # Méthodes/propriétés
        r'[;{}]\s*$',                                         # Fin de ligne avec ; ou {}
        
        # Python spécifique
        r'^\s*(import|from|class|def|if|elif|else|try|except|finally|with|for|while|return|yield|pass|break|continue)\s+',
        r'^\s*@\w+',                                          # Décorateurs
        r'^\s*#.*$',                                          # Commentaires Python
        
        # JavaScript/TypeScript
        r'^\s*(const|let|var|function|class|interface|type|enum|export|import)\s+',
        r'^\s*//.*$',                                         # Commentaires //
        r'console\.(log|error|warn|info)',
        
        # Java/C#/C++
        r'^\s*(public|private|protected|static|final|abstract|class|interface|enum|struct)\s+',
        r'^\s*[a-zA-Z_][a-zA-Z0-9_]*\s+[a-zA-Z_][a-zA-Z0-9_]*\s*[;=\(]',  # Type variable
        r'System\.(out|err)\.print',
        r'^\s*/\*.*\*/',                                      # Commentaires /* */
        
        # C/C++
        r'^\s*#(include|define|ifdef|ifndef|endif|pragma)',
        r'^\s*(int|char|float|double|void|bool|long|short|unsigned)\s+\w+',
        r'printf\s*\(',
        
        # SQL
        r'^\s*(SELECT|INSERT|UPDATE|DELETE|CREATE|DROP|ALTER|FROM|WHERE|JOIN|GROUP BY|ORDER BY|HAVING)\s+',
        
        # HTML/XML
        r'^\s*<[^>]+>.*</[^>]+>\s*$',                        # Tags complets
        r'^\s*<[^/>]+/>\s*$',                                # Tags auto-fermants
        
        # CSS
        r'^\s*[.#]?[a-zA-Z_-]+\s*\{',                       # Sélecteurs
        r'^\s*[a-zA-Z-]+\s*:\s*[^;]+;\s*$',                 # Propriétés
        
        # Shell/Bash
        r'^\s*(ls|cd|mkdir|rm|cp|mv|grep|find|cat|echo|chmod|sudo|git)\s+',
        r'^\s*[a-zA-Z_][a-zA-Z0-9_]*=\$',                   # Variables bash
        
        # Autres indicateurs
        r'^\s*\$\s*',                                        # Prompt shell
        r'>>>\s*',                                           # Prompt Python
        r'[a-zA-Z_][a-zA-Z0-9_]*::\w+',                     # Namespaces C++
        r'[a-zA-Z_]\w*\[\d+\]',                             # Arrays
        r'[a-zA-Z_]\w*\s*\*\s*\w+',                         # Pointeurs
        r'this\.|self\.',                                    # Références objet
    ]
    
    return any(re.search(pattern, text, re.IGNORECASE) for pattern in code_patterns)

def is_math_formula(text):
    """Détecte les formules mathématiques"""
    math_patterns = [
        r'[∑∏∫∂∆∇∞±≤≥≠≈∈∉∪∩∴∵α-ωΑ-Ω]',  # Symboles mathématiques
        r'\b(sin|cos|tan|log|ln|exp|sqrt|lim|max|min|sum|prod)\b',
        r'[a-zA-Z]\s*[₀-₉⁰-⁹]',                              # Indices/exposants
        r'[a-zA-Z]\^[a-zA-Z0-9]',                            # Exposants avec ^
        r'\b\d+/\d+\b',                                      # Fractions
        r'[a-zA-Z]\s*=\s*[a-zA-Z0-9+\-*/\(\)\s]+',         # Équations
    ]
    
    return any(re.search(pattern, text) for pattern in math_patterns)

def is_diagram_element(text):
    """Détecte les éléments de diagrammes/graphiques"""
    diagram_patterns = [
        # Éléments UML
        r'^[A-Z][a-zA-Z0-9]*$',                             # Noms de classes
        r'^\s*[+-]\s*\w+\s*:\s*\w+',                        # Attributs UML (+/-)
        r'^\s*[+-]\s*\w+\s*\([^)]*\)\s*:\s*\w+',           # Méthodes UML
        r'\w+\s*:\s*\w+',                                    # type:name
        
        # Éléments de graphiques
        r'\b(Figure|Fig|Graph|Chart|Diagram|Schema|Table)\s*\d+',
        r'\b(Axe|Axis|X|Y)\s*[:=]',
        r'\b(min|max|moyenne|mean|médiane|median)\s*[:=]',
        
        # Symboles de diagrammes
        r'[→←↑↓↔↕⟶⟵⟷]',                                   # Flèches
        r'[□■○●△▲◇◆]',                                       # Formes géométriques
        r'\|\s*\|',                                          # Barres parallèles
        
        # Structures de données
        r'\w+\s*->\s*\w+',                                   # Pointeurs/liens
        r'\[\s*\w+\s*\]\s*->\s*\[\s*\w+\s*\]',             # Boîtes liées
    ]
    
    return any(re.search(pattern, text) for pattern in diagram_patterns)

def is_table_like(text):
    """Détecte les tableaux"""
    return (
        "\t" in text or
        (text.count("|") >= 2 and len(text.split("|")) >= 3) or
        (text.count(",") >= 3 and len(text) < 150) or
        (re.search(r"\b\d{1,3}(\.\d{1,2})?\b", text) and "|" in text) or
        re.search(r'^[|\s]*([A-Za-z0-9\s]+\s*\|\s*){2,}', text)
    )

def detect_content_type(text):
    """Détecte le type de contenu"""
    # Ordre d'importance dans la détection
    if is_code_line(text):
        return "code_line"
    elif is_math_formula(text):
        return "formula"
    elif is_diagram_element(text):
        return "diagram_element"
    elif is_table_like(text):
        return "table_row"
    elif len(text) < 50 and text.isupper():
        return "title"
    elif text.endswith(":") and len(text) < 100:
        return "section"
    else:
        return "paragraph"

def group_similar_blocks(content):
    """Regroupe les blocs similaires consécutifs"""
    if not content:
        return content
    
    grouped = []
    current_group = [content[0]]
    current_type = content[0]["type"]
    
    for item in content[1:]:
        # Types qui peuvent être groupés
        groupable_types = ["code_line", "formula", "diagram_element", "table_row"]
        
        if (item["type"] == current_type and 
            current_type in groupable_types and
            item.get("page", 1) == current_group[-1].get("page", 1)):
            # Ajouter au groupe actuel
            current_group.append(item)
        else:
            # Finaliser le groupe précédent
            if len(current_group) > 1 and current_type in groupable_types:
                # Déterminer le type de groupe
                group_type_map = {
                    "code_line": "code_block",
                    "formula": "formula_block", 
                    "diagram_element": "diagram",
                    "table_row": "table"
                }
                
                grouped_text = "\n".join([g["text"] for g in current_group])
                grouped.append({
                    "source": current_group[0]["source"],
                    "page": current_group[0].get("page", current_group[0].get("slide", 1)),
                    "type": group_type_map[current_type],
                    "text": grouped_text
                })
            else:
                # Ajouter tel quel si pas de regroupement
                grouped.extend(current_group)
            
            # Commencer un nouveau groupe
            current_group = [item]
            current_type = item["type"]
    
    # Finaliser le dernier groupe
    if len(current_group) > 1 and current_type in ["code_line", "formula", "diagram_element", "table_row"]:
        group_type_map = {
            "code_line": "code_block",
            "formula": "formula_block", 
            "diagram_element": "diagram",
            "table_row": "table"
        }
        grouped_text = "\n".join([g["text"] for g in current_group])
        grouped.append({
            "source": current_group[0]["source"],
            "page": current_group[0].get("page", current_group[0].get("slide", 1)),
            "type": group_type_map[current_type],
            "text": grouped_text
        })
    else:
        grouped.extend(current_group)
    
    return grouped

def extract_pdf(filepath):
    doc = fitz.open(filepath)
    content = []
    
    for page_num, page in enumerate(doc, start=1):
        blocks = page.get_text("dict")["blocks"]
        
        for block in blocks:
            if "lines" not in block:
                continue
            
            text = ""
            for line in block["lines"]:
                for span in line["spans"]:
                    text += span["text"] + " "
            
            text = clean_text(text)
            if not text:
                continue
            
            block_type = detect_content_type(text)
            
            content.append({
                "source": os.path.basename(filepath),
                "page": page_num,
                "type": block_type,
                "text": text
            })
    
    return group_similar_blocks(content)

def extract_docx(filepath):
    doc = Document(filepath)
    content = []
    
    for para in doc.paragraphs:
        text = clean_text(para.text)
        if not text:
            continue
        
        style = para.style.name.lower()
        
        if "heading" in style:
            block_type = "title"
        else:
            block_type = detect_content_type(text)
        
        content.append({
            "source": os.path.basename(filepath),
            "type": block_type,
            "text": text
        })
    
    return group_similar_blocks(content)

def extract_pptx(filepath):
    prs = Presentation(filepath)
    content = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = clean_text(shape.text)
                if not text:
                    continue
                
                block_type = detect_content_type(text)
                
                content.append({
                    "source": os.path.basename(filepath),
                    "slide": slide_num,
                    "type": block_type,
                    "text": text
                })
    
    return group_similar_blocks(content)

def extract_any(filepath):
    ext = os.path.splitext(filepath)[-1].lower()
    if ext == ".pdf":
        return extract_pdf(filepath)
    elif ext == ".docx":
        return extract_docx(filepath)
    elif ext == ".pptx":
        return extract_pptx(filepath)
    else:
        raise ValueError("Format non supporté : " + ext)

# ---- MAIN ----

if __name__ == "__main__":
    INPUT_FILE = "JavaLesBases.pdf"  # ← change selon ton fichier
    output = extract_any(INPUT_FILE)

    # Sauvegarde JSON
    with open("output_structured.json", "w", encoding="utf-8") as f:
        json.dump(output, f, ensure_ascii=False, indent=2)

    print(f"✅ Extraction structurée sauvegardée dans output_structured.json ({len(output)} blocs)")

    # Sauvegarde SQLite
    import sqlite3
    db_path = "corpus.db"
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute('''
        CREATE TABLE IF NOT EXISTS blocs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            source TEXT,
            page INTEGER,
            bloc_type TEXT,
            contenu TEXT
        )
    ''')
    for bloc in output:
        c.execute('INSERT INTO blocs (source, page, bloc_type, contenu) VALUES (?, ?, ?, ?)',
                  (bloc.get('source'), bloc.get('page'), bloc.get('type'), bloc.get('text')))
    conn.commit()
    conn.close()
    print(f"✅ Extraction structurée sauvegardée aussi dans {db_path} ({len(output)} blocs)")
    
    # Afficher un aperçu des types de blocs détectés
    type_counts = {}
    for item in output:
        type_counts[item["type"]] = type_counts.get(item["type"], 0) + 1
    
    print("\n📊 Types de blocs détectés :")
    for block_type, count in sorted(type_counts.items()):
        print(f"  - {block_type}: {count} blocs")
    
    # Aperçu de quelques blocs par type
    print("\n🔍 Aperçu des contenus :")
    for block_type in ["code_block", "formula_block", "diagram", "table"]:
        examples = [item for item in output if item["type"] == block_type]
        if examples:
            print(f"\n{block_type.upper()} (exemple) :")
            print(f"  {examples[0]['text'][:100]}...")