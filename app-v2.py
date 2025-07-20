import os
import re
import json
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

def clean_text(text):
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)         # Lignes non terminées
    text = re.sub(r'\s{2,}', ' ', text)                  # Espaces multiples
    text = re.sub(r'(\w)-\s+(\w)', r'\1\2', text)         # Coupures avec "-"
    text = re.sub(r'\[\s*\]', '', text)                  # [ ] orphelins
    return text.strip()

def is_table_like(text):
    return ("\t" in text or
            (text.count(",") >= 3 and len(text) < 150) or
            (re.search(r"\b\d{1,3}(\.\d{1,2})?\b", text) and "|" in text))

def extract_pdf(filepath):
    doc = fitz.open(filepath)
    content = []
    for page_num, page in enumerate(doc, start=1):
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if "lines" not in block:
                continue
            text = ""
            for line in block["lines"]:
                for span in line["spans"]:
                    text += span["text"] + " "
            text = clean_text(text)
            if not text:
                continue
            block_type = "paragraph"
            if len(text) < 50 and text.isupper():
                block_type = "title"
            elif text.endswith(":"):
                block_type = "section"
            elif is_table_like(text):
                block_type = "table_row"
            content.append({
                "source": os.path.basename(filepath),
                "page": page_num,
                "type": block_type,
                "text": text
            })
    return content

def extract_docx(filepath):
    doc = Document(filepath)
    content = []
    for para in doc.paragraphs:
        text = clean_text(para.text)
        if not text:
            continue
        style = para.style.name.lower()
        block_type = "paragraph"
        if "heading" in style:
            block_type = "title"
        elif text.endswith(":"):
            block_type = "section"
        content.append({
            "source": os.path.basename(filepath),
            "type": block_type,
            "text": text
        })
    return content

def extract_pptx(filepath):
    prs = Presentation(filepath)
    content = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = clean_text(shape.text)
                if not text:
                    continue
                block_type = "paragraph"
                if len(text) < 50 and text.isupper():
                    block_type = "title"
                elif text.endswith(":"):
                    block_type = "section"
                content.append({
                    "source": os.path.basename(filepath),
                    "slide": slide_num,
                    "type": block_type,
                    "text": text
                })
    return content

def extract_any(filepath):
    ext = os.path.splitext(filepath)[-1].lower()
    if ext == ".pdf":
        return extract_pdf(filepath)
    elif ext == ".docx":
        return extract_docx(filepath)
    elif ext == ".pptx":
        return extract_pptx(filepath)
    else:
        raise ValueError("Format non supporté : " + ext)

# ---- MAIN ----

if __name__ == "__main__":
    INPUT_FILE = "JavaLesBases.pdf"  # ← change selon ton fichier
    output = extract_any(INPUT_FILE)

    with open("output_structured.json", "w", encoding="utf-8") as f:
        json.dump(output, f, ensure_ascii=False, indent=2)

    print(f"✅ Extraction structurée sauvegardée dans output_structured.json ({len(output)} blocs)")
